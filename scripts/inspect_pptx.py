from __future__ import annotations

from dotenv import load_dotenv

load_dotenv()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json

PPTX_PATH = "docs/poster/pptTemplate.pptx"


def emu_to_inches(emu: int) -> float:
    return round(emu / 914400, 2)


def emu_to_cm(emu: int) -> float:
    return round(emu / 360000, 2)


def rgb_to_hex(rgb: RGBColor) -> str:
    return f"#{rgb}"


def inspect_template(path: str) -> None:
    prs = Presentation(path)

    # --- Slide dimensions ---
    width_emu = prs.slide_width
    height_emu = prs.slide_height
    print("=" * 70)
    print("SLIDE DIMENSIONS")
    print("=" * 70)
    print(f"  Width:  {emu_to_inches(width_emu)} in  ({emu_to_cm(width_emu)} cm)  ({width_emu} EMU)")
    print(f"  Height: {emu_to_inches(height_emu)} in  ({emu_to_cm(height_emu)} cm)  ({height_emu} EMU)")
    print(f"  Aspect ratio: {width_emu / height_emu:.3f}")
    print()

    # --- Theme colors ---
    print("=" * 70)
    print("THEME / SLIDE MASTER COLORS")
    print("=" * 70)
    for i, master in enumerate(prs.slide_masters):
        print(f"\n  Slide Master {i}:")
        theme_elem = master.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
        )
        # Try to extract theme colors from the XML
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_scheme = master.element.find(".//a:clrScheme", ns)
        if clr_scheme is not None:
            print(f"    Color scheme name: {clr_scheme.get('name', 'unnamed')}")
            for child in clr_scheme:
                tag = child.tag.split("}")[-1]
                # Each color element has a child with val attribute
                for color_elem in child:
                    color_tag = color_elem.tag.split("}")[-1]
                    val = color_elem.get("val", color_elem.get("lastClr", ""))
                    print(f"    {tag:20s} ({color_tag}): #{val}")

        # Background
        bg = master.background
        if bg and bg.fill:
            fill = bg.fill
            print(f"\n    Background fill type: {fill.type}")
            try:
                if fill.fore_color and fill.fore_color.rgb:
                    print(f"    Background color: #{fill.fore_color.rgb}")
            except Exception:
                print("    Background color: (theme/inherited)")
    print()

    # --- Slide layouts ---
    print("=" * 70)
    print("SLIDE LAYOUTS")
    print("=" * 70)
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            print(f"\n  Layout: '{layout.name}'")
            for ph in layout.placeholders:
                left = emu_to_inches(ph.left) if ph.left else "?"
                top = emu_to_inches(ph.top) if ph.top else "?"
                w = emu_to_inches(ph.width) if ph.width else "?"
                h = emu_to_inches(ph.height) if ph.height else "?"
                print(f"    Placeholder idx={ph.placeholder_format.idx}, "
                      f"type={ph.placeholder_format.type}, "
                      f"pos=({left}, {top}), size=({w} x {h})")
    print()

    # --- Actual slides ---
    print("=" * 70)
    print(f"SLIDES ({len(prs.slides)} total)")
    print("=" * 70)

    all_colors = set()
    all_fonts = set()
    all_font_sizes = set()

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {slide_idx + 1} ---")

        # Slide background
        bg_fill = slide.background.fill
        print(f"  Background fill type: {bg_fill.type}")
        try:
            if bg_fill.fore_color and bg_fill.fore_color.rgb:
                bg_hex = f"#{bg_fill.fore_color.rgb}"
                print(f"  Background color: {bg_hex}")
                all_colors.add(("background", bg_hex))
        except Exception:
            print("  Background color: (theme/inherited)")

        for shape in slide.shapes:
            print(f"\n  Shape: '{shape.name}' (type={shape.shape_type})")
            print(f"    Position: ({emu_to_inches(shape.left)}, {emu_to_inches(shape.top)}) in")
            print(f"    Size: {emu_to_inches(shape.width)} x {emu_to_inches(shape.height)} in")

            # Shape fill
            if hasattr(shape, "fill"):
                try:
                    fill = shape.fill
                    if fill.type is not None:
                        print(f"    Fill type: {fill.type}")
                        try:
                            if fill.fore_color and fill.fore_color.rgb:
                                hex_val = f"#{fill.fore_color.rgb}"
                                print(f"    Fill color: {hex_val}")
                                all_colors.add(("shape_fill", hex_val))
                        except Exception:
                            pass
                except Exception:
                    pass

            # Shape line/border
            if hasattr(shape, "line"):
                try:
                    line = shape.line
                    if line.fill and line.fill.type is not None:
                        try:
                            if line.color and line.color.rgb:
                                hex_val = f"#{line.color.rgb}"
                                print(f"    Line color: {hex_val}")
                                all_colors.add(("line", hex_val))
                        except Exception:
                            pass
                except Exception:
                    pass

            # Text content
            if shape.has_text_frame:
                tf = shape.text_frame
                print(f"    Text frame margins: "
                      f"L={emu_to_inches(tf.margin_left or 0)}, "
                      f"R={emu_to_inches(tf.margin_right or 0)}, "
                      f"T={emu_to_inches(tf.margin_top or 0)}, "
                      f"B={emu_to_inches(tf.margin_bottom or 0)} in")
                print(f"    Word wrap: {tf.word_wrap}")

                for para_idx, para in enumerate(tf.paragraphs):
                    align = para.alignment
                    text = para.text[:80] if para.text else ""
                    if text:
                        print(f"    Paragraph {para_idx}: '{text}'")
                        if align:
                            print(f"      Alignment: {align}")

                    # Paragraph-level font
                    pf = para.font
                    if pf.size:
                        sz = pf.size.pt
                        print(f"      Para font size: {sz} pt")
                        all_font_sizes.add(sz)
                    if pf.name:
                        print(f"      Para font name: {pf.name}")
                        all_fonts.add(pf.name)

                    for run in para.runs:
                        font = run.font
                        run_text = run.text[:60] if run.text else ""
                        info_parts = []
                        if font.name:
                            info_parts.append(f"font={font.name}")
                            all_fonts.add(font.name)
                        if font.size:
                            sz = font.size.pt
                            info_parts.append(f"size={sz}pt")
                            all_font_sizes.add(sz)
                        if font.bold:
                            info_parts.append("BOLD")
                        if font.italic:
                            info_parts.append("italic")
                        if font.color and font.color.type is not None:
                            try:
                                if font.color.rgb:
                                    hex_val = f"#{font.color.rgb}"
                                    info_parts.append(f"color={hex_val}")
                                    all_colors.add(("text", hex_val))
                            except Exception:
                                try:
                                    idx = font.color.theme_color
                                    info_parts.append(f"theme_color={idx}")
                                except Exception:
                                    pass
                        if info_parts:
                            print(f"      Run: '{run_text}' [{', '.join(info_parts)}]")

    # --- Summary ---
    print("\n" + "=" * 70)
    print("SUMMARY")
    print("=" * 70)
    print(f"\n  Dimensions: {emu_to_inches(width_emu)} x {emu_to_inches(height_emu)} in")
    print(f"\n  All fonts found: {sorted(all_fonts)}")
    print(f"  All font sizes found: {sorted(all_font_sizes)}")
    print(f"\n  All colors found:")
    for context, color in sorted(all_colors):
        print(f"    {context:20s} {color}")


def inspect_xml_details(path: str) -> None:
    """Dig into the raw XML for theme colors, default fonts, and shape colors."""
    from lxml import etree

    prs = Presentation(path)

    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    print("\n" + "=" * 70)
    print("DETAILED XML INSPECTION")
    print("=" * 70)

    # --- Theme element from slide master ---
    for i, master in enumerate(prs.slide_masters):
        # The theme is linked via relationship
        theme_part = None
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break

        if theme_part:
            theme_xml = etree.fromstring(theme_part.blob)
            print(f"\n--- Theme from Slide Master {i} ---")

            # Color scheme
            clr_scheme = theme_xml.find(".//a:clrScheme", ns)
            if clr_scheme is not None:
                print(f"\n  Color Scheme: '{clr_scheme.get('name', 'unnamed')}'")
                for child in clr_scheme:
                    tag = child.tag.split("}")[-1]
                    for color_elem in child:
                        color_tag = color_elem.tag.split("}")[-1]
                        val = color_elem.get("val", color_elem.get("lastClr", ""))
                        print(f"    {tag:20s} -> #{val}")

            # Font scheme
            font_scheme = theme_xml.find(".//a:fontScheme", ns)
            if font_scheme is not None:
                print(f"\n  Font Scheme: '{font_scheme.get('name', 'unnamed')}'")
                major = font_scheme.find("a:majorFont", ns)
                minor = font_scheme.find("a:minorFont", ns)
                if major is not None:
                    latin = major.find("a:latin", ns)
                    if latin is not None:
                        print(f"    Major (headings): {latin.get('typeface', '?')}")
                if minor is not None:
                    latin = minor.find("a:latin", ns)
                    if latin is not None:
                        print(f"    Minor (body):     {latin.get('typeface', '?')}")

            # Format scheme (background fills etc)
            fmt_scheme = theme_xml.find(".//a:fmtScheme", ns)
            if fmt_scheme is not None:
                print(f"\n  Format Scheme: '{fmt_scheme.get('name', 'unnamed')}'")

    # --- Per-shape detailed XML for colors ---
    print("\n--- Shape fill/text colors from XML ---")
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            elem = shape.element
            name = shape.name

            # Solid fills on shapes
            solid_fills = elem.findall(".//a:solidFill", ns)
            for sf in solid_fills:
                for child in sf:
                    tag = child.tag.split("}")[-1]
                    val = child.get("val", "")
                    if tag == "srgbClr":
                        print(f"  {name}: solidFill srgbClr #{val}")
                    elif tag == "schemeClr":
                        print(f"  {name}: solidFill schemeClr {val}")

            # Gradient fills
            grad_fills = elem.findall(".//a:gradFill", ns)
            for gf in grad_fills:
                stops = gf.findall(".//a:gs", ns)
                for stop in stops:
                    pos = stop.get("pos", "?")
                    for child in stop:
                        for cc in child:
                            tag = cc.tag.split("}")[-1]
                            if tag == "srgbClr":
                                print(f"  {name}: gradFill pos={pos} #{cc.get('val', '')}")
                            elif tag == "schemeClr":
                                print(f"  {name}: gradFill pos={pos} scheme={cc.get('val', '')}")
                        # Direct color on the stop child
                        ctag = child.tag.split("}")[-1]
                        val = child.get("val", "")
                        if ctag == "srgbClr":
                            print(f"  {name}: gradFill pos={pos} #{val}")
                        elif ctag == "schemeClr":
                            print(f"  {name}: gradFill pos={pos} scheme={val}")

            # Text run font sizes from XML (catches defaults we might miss)
            rPrs = elem.findall(".//a:rPr", ns)
            for rPr in rPrs:
                sz = rPr.get("sz")
                b = rPr.get("b")
                i_attr = rPr.get("i")
                latin = rPr.find("a:latin", ns)
                font_name = latin.get("typeface") if latin is not None else None

                solid = rPr.find("a:solidFill", ns)
                color_str = ""
                if solid is not None:
                    for c in solid:
                        ctag = c.tag.split("}")[-1]
                        if ctag == "srgbClr":
                            color_str = f" color=#{c.get('val', '')}"
                        elif ctag == "schemeClr":
                            color_str = f" scheme={c.get('val', '')}"

                parts = []
                if font_name:
                    parts.append(f"font={font_name}")
                if sz:
                    parts.append(f"size={int(sz)/100}pt")
                if b == "1":
                    parts.append("BOLD")
                if i_attr == "1":
                    parts.append("italic")
                if color_str:
                    parts.append(color_str.strip())

                if parts:
                    # Get some text context
                    parent_r = rPr.getparent()
                    t_elem = parent_r.find("a:t", ns) if parent_r is not None else None
                    text_preview = (t_elem.text or "")[:40] if t_elem is not None else ""
                    print(f"  {name}: rPr [{', '.join(parts)}] text='{text_preview}'")

            # Default text (defRPr)
            defRPrs = elem.findall(".//a:defRPr", ns)
            for drp in defRPrs:
                sz = drp.get("sz")
                latin = drp.find("a:latin", ns)
                font_name = latin.get("typeface") if latin is not None else None
                solid = drp.find("a:solidFill", ns)
                color_str = ""
                if solid is not None:
                    for c in solid:
                        ctag = c.tag.split("}")[-1]
                        if ctag == "srgbClr":
                            color_str = f" color=#{c.get('val', '')}"
                        elif ctag == "schemeClr":
                            color_str = f" scheme={c.get('val', '')}"

                parts = []
                if font_name:
                    parts.append(f"font={font_name}")
                if sz:
                    parts.append(f"size={int(sz)/100}pt")
                if color_str:
                    parts.append(color_str.strip())
                if parts:
                    print(f"  {name}: defRPr [{', '.join(parts)}]")


if __name__ == "__main__":
    inspect_template(PPTX_PATH)
    inspect_xml_details(PPTX_PATH)
